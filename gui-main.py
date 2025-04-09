import os
import fitz
import sys
import fnmatch
import logging
from docx import Document  # For .docx files
from pptx import Presentation  # For .pptx files
from openpyxl import load_workbook  # For .xlsx files
import chardet
import customtkinter as ctk
from tkinter import filedialog, messagebox
import threading
import time

# Configure logging for successful file paths
success_logger = logging.getLogger("success_logger")
success_logger.setLevel(logging.INFO)
success_handler = logging.FileHandler("success.log")
success_handler.setFormatter(logging.Formatter("%(asctime)s - %(message)s"))
success_logger.addHandler(success_handler)

# Configure logging for unsuccessful file paths
unsuccessful_logger = logging.getLogger("unsuccessful_logger")
unsuccessful_logger.setLevel(logging.WARNING)
unsuccessful_handler = logging.FileHandler("unsuccessful.log")
unsuccessful_handler.setFormatter(logging.Formatter("%(asctime)s - %(message)s"))
unsuccessful_logger.addHandler(unsuccessful_handler)

# Store handlers in a list for later access
log_handlers = [success_handler, unsuccessful_handler]

def clear_log_files():
    log_files = ["success.log", "unsuccessful.log"]
    
    # Close all logging handlers
    for handler in log_handlers:
        handler.close()
    
    # Delete log files
    for log_file in log_files:
        if os.path.exists(log_file):
            os.remove(log_file)

# Extract text from .pdf file
def extract_text_from_pdf(pdf_path):
    try:
        text = ""
        doc = fitz.open(pdf_path)

        for page in doc:
            text += page.get_text("text") + "\n"

        success_logger.info(f"Successfully processed: {pdf_path}")
        return text.strip()
    except Exception as e:
        error_message = f"Error extracting text from PDF {pdf_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {pdf_path} - {error_message}")
        return ""

# Extract text from .txt file
def extract_text_from_txt(file_path):
    try:
        # Detect file encoding
        with open(file_path, "rb") as f:
            raw_data = f.read(10000)  # Read a sample of the file to detect encoding
            detected = chardet.detect(raw_data)
            encoding = detected["encoding"] if detected["encoding"] else "utf-8"

        # Read file using detected encoding
        with open(file_path, "r", encoding=encoding, errors="replace") as file:
            text = file.read().strip()
            success_logger.info(f"Successfully processed: {file_path}")
            return text
    except Exception as e:
        error_message = f"Error reading .txt file {file_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {file_path} - {error_message}")
        return ""

# Extract text from .docx file
def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        success_logger.info(f"Successfully processed: {file_path}")
        return text.strip()
    except Exception as e:
        error_message = f"Error reading .docx file {file_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {file_path} - {error_message}")
        return ""

# Extract text from .pptx file
def extract_text_from_pptx(file_path):
    try:
        ppt = Presentation(file_path)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        success_logger.info(f"Successfully processed: {file_path}")
        return text.strip()
    except Exception as e:
        error_message = f"Error reading .pptx file {file_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {file_path} - {error_message}")
        return ""

# Extract text from .xlsx file
def extract_text_from_xlsx(file_path):
    try:
        workbook = load_workbook(file_path)
        text = ""
        for sheet in workbook:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell)
                text += row_text + "\n"
        success_logger.info(f"Successfully processed: {file_path}")
        return text.strip()
    except Exception as e:
        error_message = f"Error reading .xlsx file {file_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {file_path} - {error_message}")
        return ""

# Search for text in a file
def search_text_in_file(file_path, search_text):
    try:
        if file_path.endswith(".pdf"):
            text = extract_text_from_pdf(file_path)
        elif file_path.endswith(".txt"):
            text = extract_text_from_txt(file_path)
        elif file_path.endswith(".docx"):
            text = extract_text_from_docx(file_path)
        elif file_path.endswith(".pptx"):
            text = extract_text_from_pptx(file_path)
        elif file_path.endswith(".xlsx"):
            text = extract_text_from_xlsx(file_path)
        else:
            return False  # Skip unsupported file types

        # Check for exact case-sensitive match
        if search_text in text:
            return True
        return False
    except Exception as e:
        error_message = f"Error searching in file {file_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {file_path} - {error_message}")
        return False

# Recursively search for files and check for text
def search_files(directory, search_text, file_extension=None):
    matches = []
    for root, _, files in os.walk(directory, onerror=lambda e: None):
        for file in files:
            file_path = os.path.join(root, file)
            # Check if the file matches the extension (if provided)
            if file_extension and not fnmatch.fnmatch(file, f"*{file_extension}"):
                continue
            if search_text_in_file(file_path, search_text):
                matches.append(file_path)  # Add matching file to the list
    return matches

# Get root directories based on OS
def get_root_directories():
    if sys.platform == "linux":
        return ["/"]
    elif sys.platform == "win32":
        return [
            f"{chr(65 + i)}:/" for i in range(26) if os.path.exists(f"{chr(65 + i)}:/")
        ]
    else:
        raise Exception("Unsupported operating system")

# GUI Application
class SearchApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("File Search Tool")
        self.geometry("600x450")
        self.configure(bg="#2E3440")

        # Theme and appearance
        ctk.set_appearance_mode("Dark")
        ctk.set_default_color_theme("blue")

        # Input Frame
        self.input_frame = ctk.CTkFrame(self, fg_color="#3B4252")
        self.input_frame.pack(pady=20, padx=20, fill="x")

        # Enter Text to Search
        self.search_label = ctk.CTkLabel(self.input_frame, text="Enter Text to Search:", font=("Arial", 14))
        self.search_label.grid(row=0, column=0, padx=5, pady=5, sticky="w")

        self.search_entry = ctk.CTkEntry(self.input_frame, width=300, font=("Arial", 12))
        self.search_entry.grid(row=0, column=1, padx=5, pady=5, sticky="ew")

        # Select Directory
        self.directory_label = ctk.CTkLabel(self.input_frame, text="Select Directory:", font=("Arial", 14))
        self.directory_label.grid(row=1, column=0, padx=5, pady=5, sticky="w")

        self.directory_var = ctk.StringVar(value="All Directories")
        self.directory_dropdown = ctk.CTkOptionMenu(
            self.input_frame,
            values=["All Directories"] + get_root_directories(),
            variable=self.directory_var,
            font=("Arial", 12),
        )
        self.directory_dropdown.grid(row=1, column=1, padx=5, pady=5, sticky="ew")

        # Filter by File Type
        self.file_type_label = ctk.CTkLabel(self.input_frame, text="Filter by File Type:", font=("Arial", 14))
        self.file_type_label.grid(row=2, column=0, padx=5, pady=5, sticky="w")

        self.file_type_var = ctk.StringVar(value="All Files")
        self.file_type_dropdown = ctk.CTkOptionMenu(
            self.input_frame,
            values=["All Files", ".txt", ".pdf", ".docx", ".pptx", ".xlsx"],
            variable=self.file_type_var,
            font=("Arial", 12),
        )
        self.file_type_dropdown.grid(row=2, column=1, padx=5, pady=5, sticky="ew")

        # Search Button
        self.search_button = ctk.CTkButton(
            self.input_frame, text="Search", command=self.start_search, font=("Arial", 14), fg_color="#5E81AC"
        )
        self.search_button.grid(row=3, column=0, columnspan=2, pady=10)

        # Output Frame
        self.output_frame = ctk.CTkFrame(self, fg_color="#3B4252")
        self.output_frame.pack(pady=20, padx=20, fill="x")

        self.output_label = ctk.CTkLabel(self.output_frame, text="Result:", font=("Arial", 14))
        self.output_label.pack(pady=5)

        # Scrollable Textbox for Results
        self.output_text = ctk.CTkTextbox(self.output_frame, width=550, height=150, font=("Arial", 12), wrap="none")
        self.output_text.pack(pady=5, fill="both", expand=True)

        # Add a scrollbar to the textbox
        self.output_text.configure(state="normal")  # Allow editing for inserting results

        # Loading Animation Label
        self.loading_label = ctk.CTkLabel(self.output_frame, text="", font=("Arial", 14))
        self.loading_label.pack(pady=5)

        # Clear Logs Button
        self.clear_logs_button = ctk.CTkButton(
            self, text="Clear Logs", command=clear_log_files, font=("Arial", 14), fg_color="#BF616A"
        )
        self.clear_logs_button.pack(pady=10)

        # Store the matched file paths
        self.matched_file_paths = []

    def start_search(self):
        search_text = self.search_entry.get().strip()
        if not search_text or search_text.isspace():
            messagebox.showwarning("Input Error", "Please enter valid text to search.")
            return

        file_extension = self.file_type_var.get()
        if file_extension == "All Files":
            file_extension = None

        selected_directory = self.directory_var.get()
        if selected_directory == "All Directories":
            search_paths = get_root_directories()
        else:
            search_paths = [selected_directory]

        clear_log_files()  # Clear logs before starting a new search

        self.output_text.delete("1.0", "end")  # Clear previous output

        # Start loading animation
        self.loading_label.configure(text="Searching...")
        self.update_loading_animation()

        # Run search in a separate thread to avoid freezing the UI
        threading.Thread(target=self.perform_search, args=(search_paths, search_text, file_extension), daemon=True).start()

    def perform_search(self, search_paths, search_text, file_extension):
        matches = []
        for path in search_paths:
            matches.extend(search_files(path, search_text, file_extension))

        # Stop loading animation and display results
        self.after(0, self.display_results, matches)

    def display_results(self, matches):
        self.loading_label.configure(text="")  # Stop loading animation
        self.output_text.delete("1.0", "end")  # Clear previous output

        if matches:
            self.output_text.insert("end", f"Found {len(matches)} matches:\n\n")
            for match in matches:
                # Insert each match as a clickable link in the textbox
                self.output_text.insert("end", f"{match}\n")
        else:
            self.output_text.insert("end", "No matches found.")

    def update_loading_animation(self):
        current_text = self.loading_label.cget("text")
        if "Searching" in current_text:
            if current_text.endswith("..."):
                self.loading_label.configure(text="Searching")
            else:
                self.loading_label.configure(text=current_text + ".")
            self.after(500, self.update_loading_animation)  # Update every 500ms

# Run the application
if __name__ == "__main__":
    app = SearchApp()
    app.mainloop()