import os
import fitz  # PyMuPDF
import sys
import fnmatch
import logging
from logging.handlers import RotatingFileHandler
from docx import Document  # For .docx files
from pptx import Presentation  # For .pptx files
from openpyxl import load_workbook  # For .xlsx files
import chardet
import contextlib

# Configure logging for successful file paths
success_logger = logging.getLogger("success_logger")
success_logger.setLevel(logging.INFO)
success_handler = RotatingFileHandler("success.log", encoding="utf-8")
success_handler.setFormatter(logging.Formatter("%(asctime)s - %(message)s"))
success_logger.addHandler(success_handler)

# Configure logging for unsuccessful file paths
unsuccessful_logger = logging.getLogger("unsuccessful_logger")
unsuccessful_logger.setLevel(logging.WARNING)
unsuccessful_handler = RotatingFileHandler("unsuccessful.log", encoding="utf-8")
unsuccessful_handler.setFormatter(logging.Formatter("%(asctime)s - %(message)s"))
unsuccessful_logger.addHandler(unsuccessful_handler)

# Store handlers in a list for later access
log_handlers = [success_handler, unsuccessful_handler]

def clear_log_files():
    log_files = ["success.log", "unsuccessful.log"]
    
    # Close all logging handlers
    for handler in log_handlers:
        handler.close()
    
    # Delete log files
    for log_file in log_files:
        if os.path.exists(log_file):
            os.remove(log_file)

# Extract text from .pdf file
def extract_text_from_pdf(pdf_path):
    try:
        text = ""
        # Redirect stderr to capture MuPDF errors
        with contextlib.redirect_stderr(sys.stdout):  # Temporarily redirect stderr
            doc = fitz.open(pdf_path)

            for page in doc:
                try:
                    text += page.get_text("text") + "\n"
                except Exception as page_error:
                    error_message = f"Error extracting text from page in PDF {pdf_path}: {page_error}"
                    unsuccessful_logger.warning(f"Unsuccessful: {pdf_path} - {error_message}")

        success_logger.info(f"Successfully processed: {pdf_path}")
        return text.strip()
    except Exception as e:
        error_message = f"Error opening or processing PDF {pdf_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {pdf_path} - {error_message}")
        return ""

# Extract text from .txt file
def extract_text_from_txt(file_path):
    try:
        # Detect file encoding
        with open(file_path, "rb") as f:
            raw_data = f.read(1000)  # Read a sample of the file to detect encoding
            detected = chardet.detect(raw_data)
            encoding = detected["encoding"] if detected["encoding"] else "utf-8"

        # Read file using detected encoding
        with open(file_path, "r", encoding=encoding, errors="replace") as file:
            text = file.read().strip()
            success_logger.info(f"Successfully processed: {file_path}")
            return text
    except Exception as e:
        error_message = f"Error reading .txt file {file_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {file_path} - {error_message}")
        return ""

# Extract text from .docx file
def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        success_logger.info(f"Successfully processed: {file_path}")
        return text.strip()
    except Exception as e:
        error_message = f"Error reading .docx file {file_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {file_path} - {error_message}")
        return ""

# Extract text from .pptx file
def extract_text_from_pptx(file_path):
    try:
        ppt = Presentation(file_path)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        success_logger.info(f"Successfully processed: {file_path}")
        return text.strip()
    except Exception as e:
        error_message = f"Error reading .pptx file {file_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {file_path} - {error_message}")
        return ""

# Extract text from .xlsx file
def extract_text_from_xlsx(file_path):
    try:
        workbook = load_workbook(file_path)
        text = ""
        for sheet in workbook:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell)
                text += row_text + "\n"
        success_logger.info(f"Successfully processed: {file_path}")
        return text.strip()
    except Exception as e:
        error_message = f"Error reading .xlsx file {file_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {file_path} - {error_message}")
        return ""

# Search for text in a file
def search_text_in_file(file_path, search_text):
    try:
        if file_path.endswith(".pdf"):
            text = extract_text_from_pdf(file_path)
        elif file_path.endswith(".txt"):
            text = extract_text_from_txt(file_path)
        elif file_path.endswith(".docx"):
            text = extract_text_from_docx(file_path)
        elif file_path.endswith(".pptx"):
            text = extract_text_from_pptx(file_path)
        elif file_path.endswith(".xlsx"):
            text = extract_text_from_xlsx(file_path)
        else:
            return False  # Skip unsupported file types

        # Check for exact case-sensitive match
        if search_text in text:
            return True
        return False
    except Exception as e:
        error_message = f"Error searching in file {file_path}: {e}"
        unsuccessful_logger.warning(f"Unsuccessful: {file_path} - {error_message}")
        return False

# Recursively search for files and check for text
def search_files(directory, search_text, file_extension=None):
    matches = []
    for root, _, files in os.walk(directory, onerror=lambda e: None):
        for file in files:
            try:
                file_path = os.path.join(root, file)
                # Check if the file matches the extension (if provided)
                if file_extension and not fnmatch.fnmatch(file, f"*{file_extension}"):
                    continue
                if search_text_in_file(file_path, search_text):
                    matches.append(file_path)
            except Exception as e:
                error_message = f"Error processing file {file_path}: {e}"
                unsuccessful_logger.warning(f"Unsuccessful: {file_path} - {error_message}")
    return matches

# Get root directories based on OS
def get_root_directories():
    if sys.platform == "linux":
        return ["/"]
    elif sys.platform == "win32":
        return [
            f"{chr(65 + i)}:/" for i in range(26) if os.path.exists(f"{chr(65 + i)}:/")
        ]
    else:
        raise Exception("Unsupported operating system")

# Main function
def main():
    clear_log_files()
    search_text = input("Enter the text to search: ").strip()
    if not search_text:
        print("No search text provided.")
        return

    search_paths = (
        input(
            "Enter starting path (comma-separated, e.g., C:/, E:/Games, leave blank for system-wide): "
        )
        .strip()
        .split(",")
    )
    search_paths = [path.strip() for path in search_paths if path.strip()]
    if not search_paths:
        search_paths = get_root_directories()

    file_extension = (
        input(
            "Filter by file extension (e.g., .txt, .pdf, .docx, .pptx, .xlsx, leave blank for all types): "
        ).strip()
        or None
    )

    print("\nSearching... This may take some time.")
    matches = []
    for path in search_paths:
        matches.extend(search_files(path, search_text, file_extension))

    if matches:
        print("\nMatches found in the following files:")
        for match in matches:
            print(match)
    else:
        print("\nNo matches found.")

if __name__ == "__main__":
    main()